"""
Report Service — generates PDF, PPTX, Excel reports from analysis results.
"""
import io, uuid
from pathlib import Path
from datetime import datetime
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from app.models.report import Report, ReportFormat, ReportStatus
from app.models.analysis import AnalysisResult
from app.core.config import settings
from app.core.celery_app import celery_app
from loguru import logger


class ReportService:

    async def generate(self, analysis_id: str, format: str, db: AsyncSession) -> Report:
        report = Report(
            analysis_id = analysis_id,
            format      = ReportFormat(format),
            status      = ReportStatus.generating,
        )
        db.add(report)
        await db.flush()
        generate_report_task.delay(str(report.id), analysis_id, format)
        return report


@celery_app.task(name="generate_report_task")
def generate_report_task(report_id: str, analysis_id: str, format: str):
    from app.core.database import AsyncSessionLocal
    import asyncio

    async def _run():
        async with AsyncSessionLocal() as db:
            r = await db.execute(select(Report).where(Report.id == report_id))
            report: Report = r.scalar_one()
            a = await db.execute(select(AnalysisResult).where(AnalysisResult.id == analysis_id))
            analysis: AnalysisResult = a.scalar_one()

            try:
                out_dir = Path(settings.LOCAL_UPLOAD_DIR) / "reports"
                out_dir.mkdir(parents=True, exist_ok=True)
                out_path = out_dir / f"{report_id}.{format}"

                if format == "pdf":
                    _generate_pdf(analysis, out_path)
                elif format == "pptx":
                    _generate_pptx(analysis, out_path)
                elif format == "excel":
                    _generate_excel(analysis, out_path)
                elif format == "json":
                    _generate_json(analysis, out_path)

                report.file_path    = str(out_path)
                report.download_url = f"/api/v1/reports/{report_id}/download"
                report.status       = ReportStatus.ready
                await db.commit()
                logger.info(f"Report {report_id} generated: {format}")

            except Exception as e:
                report.status = ReportStatus.error
                await db.commit()
                logger.error(f"Report {report_id} failed: {e}")

    asyncio.get_event_loop().run_until_complete(_run())


def _generate_pdf(analysis: AnalysisResult, path: Path):
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    doc    = SimpleDocTemplate(str(path), pagesize=letter)
    styles = getSampleStyleSheet()
    story  = []

    story.append(Paragraph("AI Business Data Analysis Report", styles["Title"]))
    story.append(Paragraph(f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}", styles["Normal"]))
    story.append(Spacer(1, 12))

    if analysis.business_problem:
        story.append(Paragraph("Business Problem", styles["Heading1"]))
        story.append(Paragraph(analysis.business_problem, styles["Normal"]))
        story.append(Spacer(1, 8))

    story.append(Paragraph("Key Insights", styles["Heading1"]))
    for ins in (analysis.insights or []):
        story.append(Paragraph(f"[{ins.get('severity','').upper()}] {ins.get('title','')}", styles["Heading3"]))
        story.append(Paragraph(ins.get("detail",""), styles["Normal"]))
        story.append(Spacer(1, 6))

    story.append(Paragraph("Recommendations", styles["Heading1"]))
    for rec in (analysis.recommendations or []):
        story.append(Paragraph(f"• {rec.get('title','')} — {rec.get('description','')}", styles["Normal"]))

    doc.build(story)


def _generate_pptx(analysis: AnalysisResult, path: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs   = Presentation()
    blank = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text        = "AI Business Data Analysis"
    slide.placeholders[1].text     = f"Generated {datetime.utcnow().strftime('%Y-%m-%d')}"

    # Insights slides
    for ins in (analysis.insights or [])[:5]:
        slide = prs.slides.add_slide(blank)
        slide.shapes.title.text    = ins.get("title", "Insight")
        slide.placeholders[1].text = ins.get("detail","") + "\n\nAction: " + ins.get("action","")

    prs.save(str(path))


def _generate_excel(analysis: AnalysisResult, path: Path):
    import openpyxl
    wb = openpyxl.Workbook()

    ws = wb.active
    ws.title = "Insights"
    ws.append(["Severity","Title","Detail","Action"])
    for ins in (analysis.insights or []):
        ws.append([ins.get("severity",""), ins.get("title",""), ins.get("detail",""), ins.get("action","")])

    ws2 = wb.create_sheet("KPIs")
    ws2.append(["Metric","Value","Change %","Trend"])
    for kpi in (analysis.kpis or []):
        ws2.append([kpi.get("label",""), kpi.get("value",""), kpi.get("change",""), kpi.get("trend","")])

    ws3 = wb.create_sheet("SQL Queries")
    ws3.append(["Title","Description","SQL"])
    for q in (analysis.sql_queries or []):
        ws3.append([q.get("title",""), q.get("description",""), q.get("sql","")])

    wb.save(str(path))


def _generate_json(analysis: AnalysisResult, path: Path):
    import json
    payload = {
        "id":              analysis.id,
        "generated_at":   datetime.utcnow().isoformat(),
        "business_problem":analysis.business_problem,
        "insights":        analysis.insights,
        "kpis":            analysis.kpis,
        "sql_queries":     analysis.sql_queries,
        "recommendations": analysis.recommendations,
        "anomalies":       analysis.anomalies,
    }
    path.write_text(json.dumps(payload, indent=2, default=str))


report_service = ReportService()
